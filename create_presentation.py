#!/usr/bin/env python3
"""
Script to create a PowerPoint presentation for the Decentralized Freelance Marketplace
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR

def create_hackathon_presentation():
    """Create a comprehensive PowerPoint presentation for the hackathon"""
    
    # Create presentation
    prs = Presentation()
    
    # Define colors (Orange theme matching the UI)
    ORANGE_RGB = RGBColor(255, 165, 0)  # Orange
    DARK_RGB = RGBColor(31, 41, 55)     # Dark gray
    WHITE_RGB = RGBColor(255, 255, 255) # White
    
    def add_title_slide():
        """Slide 1: Title Slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "🚀 Decentralized Freelance Marketplace"
        subtitle.text = "Revolutionizing Remote Work with Blockchain Technology\n\nBuilt on Stacks Blockchain"
        
        # Style title
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = ORANGE_RGB
        
    def add_problem_slide():
        """Slide 2: The Problem"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "The Problem 🚨"
        content.text = """Current Freelance Platforms Issues:
        
• High Fees (15-20% platform commission)
• Payment Disputes without transparent resolution  
• Centralized Control over funds and data
• Geographic Restrictions and payment barriers
• Lack of Trust between clients and freelancers
• No Skill Verification system

"$400 billion freelance economy needs decentralization" """
        
    def add_solution_slide():
        """Slide 3: Our Solution"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Our Solution 💡"
        content.text = """Decentralized Freelance Marketplace
Trustless • Transparent • Global

✅ Zero Platform Fees - Direct peer-to-peer payments
✅ Smart Contract Escrow - Automated payment security  
✅ Blockchain Transparency - All transactions visible
✅ Global Accessibility - No geographic restrictions
✅ Dispute Resolution - Community-driven arbitration
✅ Skill Verification - On-chain credential system"""
        
    def add_features_slide():
        """Slide 4: Key Features"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Key Features 🌟"
        content.text = """🔒 Smart Contract Escrow
• Funds locked until project completion
• Client approval required for payment release
• Automatic dispute resolution mechanism

👥 Verified Freelancer Network  
• Pre-vetted developers across multiple stacks
• On-chain skill verification
• Rating and review system

📊 Modern Dashboard UI
• Staking-inspired design with smooth animations
• Real-time project tracking  
• Portfolio analytics"""
        
    def add_tech_architecture_slide():
        """Slide 5: Technical Architecture"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Technical Architecture 🏗️"
        content.text = """Frontend Stack
• React 18 + TypeScript
• Framer Motion for animations
• Tailwind CSS for styling  
• Lucide Icons for UI components

Blockchain Integration
• Stacks Blockchain for smart contracts
• @stacks/connect for wallet integration
• Clarity smart contracts (planned)

Key Components
• Wallet Connection System
• Project Creator & Manager
• Freelancer Discovery
• Payment Escrow System"""
        
    def add_smart_contract_flow_slide():
        """Slide 6: Smart Contract Flow"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Smart Contract Flow 🔄"
        content.text = """1. Client Creates Project
   ↓
2. Funds Locked in Escrow
   ↓  
3. Freelancer Assigned
   ↓
4. Work Completed
   ↓
5. Client Reviews & Approves
   ↓
6. Payment Released Automatically

Security Features:
• Multi-signature approvals
• Time-locked releases
• Dispute arbitration
• Emergency withdrawals"""
        
    def add_ux_demo_slide():
        """Slide 7: User Experience Demo"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "User Experience Demo 🎨"
        content.text = """Modern Staking Dashboard Design
• Orange-themed gradient UI
• Smooth micro-interactions
• Responsive design
• Dark mode optimized

Key Screens:
1. Dashboard - Portfolio overview with stats
2. Project Creator - Intuitive project posting  
3. Freelancer Discovery - Skill-based matching
4. Payment Management - Escrow tracking

[Include screenshots or video demo here]"""
        
    def add_freelancers_slide():
        """Slide 8: Verified Freelancers"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Verified Freelancers 👨‍💻"
        content.text = """Pre-Vetted Talent Pool

Frontend Specialists:
• React/Vue.js experts
• UI/UX designers  
• Mobile developers

Backend Engineers:
• Node.js/Python developers
• Blockchain specialists
• DevOps engineers

Full-Stack Developers:
• End-to-end project delivery
• Multi-technology expertise

All with verified skills and portfolio reviews"""
        
    def add_market_opportunity_slide():
        """Slide 9: Market Opportunity"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Market Opportunity 📈"
        content.text = """Global Freelance Market
• $400B total market size (2024)
• 57 million freelancers in US alone
• 36% of workforce freelancing by 2027

Blockchain Adoption  
• $67B DeFi market cap
• Growing demand for decentralized solutions
• Increasing crypto payment adoption

Our Target
• Initial: 1,000 projects in first year
• Goal: 10% cost reduction for clients
• Revenue: Transaction volume growth"""
        
    def add_technical_achievements_slide():
        """Slide 10: Technical Achievements"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Technical Achievements 🏆"
        content.text = """What We Built:

✅ Complete Frontend Application  
• 8 major components with animations
• Wallet integration system
• Responsive design

✅ Smart Contract Integration
• Escrow payment system
• Project management functions  
• Dispute resolution framework

✅ Modern UI/UX
• Framer Motion animations
• Dashboard-style interface
• Mobile-responsive design

✅ GitHub Repository
• Clean, documented codebase
• TypeScript implementation
• Ready for deployment"""
        
    def add_business_model_slide():
        """Slide 12: Business Model"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Business Model 💰"
        content.text = """Revenue Streams

🔸 Transaction Fees (1-2% vs 15-20% traditional)
🔸 Premium Freelancer Listings
🔸 Enterprise Solutions  
🔸 Staking Rewards Distribution
🔸 Educational Courses & Certifications

Value Proposition
• For Clients: Lower costs, secure payments
• For Freelancers: Higher earnings, global reach
• For Platform: Sustainable, low-fee model"""
        
    def add_roadmap_slide():
        """Slide 13: Roadmap"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Roadmap 🗺️"
        content.text = """Phase 1: Foundation ✅
• Core UI/UX development
• Wallet integration
• Basic project management

Phase 2: Smart Contracts (Next 2 months)
• Deploy Clarity contracts
• Implement escrow system
• Add dispute resolution

Phase 3: Scale (3-6 months)  
• Mobile application
• Advanced matching algorithms
• Multi-blockchain support

Phase 4: Ecosystem (6-12 months)
• DAO governance
• Token economics
• Enterprise partnerships"""
        
    def add_call_to_action_slide():
        """Slide 18: Call to Action"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Call to Action 🎯"
        content.text = """What We're Looking For

🤝 Partnerships
• Blockchain development teams
• Freelance talent networks  
• Enterprise clients

💰 Investment
• Seed funding for smart contract development
• Marketing and user acquisition
• Technical team expansion

🚀 Community
• Early adopter freelancers
• Beta testing clients
• Developer contributors"""
        
    def add_thank_you_slide():
        """Slide 19: Thank You"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Thank You! 🙏"
        content.text = """Decentralized Freelance Marketplace
Building the Future of Work

Contact Information:
🌐 GitHub: https://github.com/surf3rr/decentralized-updated
📧 Email: [your-email@example.com]
🐦 Twitter: [@your_handle]  
💼 LinkedIn: [Your LinkedIn]

Questions & Demo
Ready for live demonstration"""
        
    # Add all slides
    add_title_slide()
    add_problem_slide()
    add_solution_slide()
    add_features_slide()
    add_tech_architecture_slide()
    add_smart_contract_flow_slide()
    add_ux_demo_slide()
    add_freelancers_slide()
    add_market_opportunity_slide()
    add_technical_achievements_slide()
    add_business_model_slide()
    add_roadmap_slide()
    add_call_to_action_slide()
    add_thank_you_slide()
    
    # Apply consistent styling to all slides
    for slide in prs.slides:
        # Style title
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            title_shape = slide.shapes.title
            if title_shape.has_text_frame:
                title_shape.text_frame.paragraphs[0].font.bold = True
                title_shape.text_frame.paragraphs[0].font.size = Pt(36)
                title_shape.text_frame.paragraphs[0].font.color.rgb = ORANGE_RGB
        
        # Style content
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.size = Pt(18)
                    paragraph.font.color.rgb = DARK_RGB
    
    return prs

def main():
    """Create and save the presentation"""
    try:
        print("Creating PowerPoint presentation...")
        prs = create_hackathon_presentation()
        
        # Save presentation
        output_file = "Decentralized_Freelance_Marketplace_Hackathon.pptx"
        prs.save(output_file)
        print(f"✅ Presentation saved as: {output_file}")
        
        print("\n🎯 Presentation Features:")
        print("• 14 professional slides")
        print("• Orange-themed design matching your UI")
        print("• Comprehensive project overview")
        print("• Technical achievements highlighted")
        print("• Market opportunity analysis")
        print("• Clear call to action")
        
        print("\n📋 Next Steps:")
        print("1. Review and customize the presentation")
        print("2. Add screenshots from your application")
        print("3. Practice your demo")
        print("4. Prepare for Q&A session")
        
    except ImportError:
        print("❌ Error: python-pptx not installed")
        print("Please run: pip install python-pptx")
        print("Then run this script again")
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")

if __name__ == "__main__":
    main()
